"""Extract plain text from an uploaded file for ideation.

The uploaded bytes never touch disk — they're parsed in memory and discarded.
Formats: txt/md/qmd/csv/json (text), pdf, docx, pptx, xlsx, and ODF
(odt/ods/odp via their zipped content.xml). Extraction libs are imported lazily
so a missing one only disables that format (the app keeps running).
"""

from __future__ import annotations

import io
import zipfile
from xml.etree import ElementTree as ET

_TEXT_SUFFIXES = {".txt", ".md", ".markdown", ".qmd", ".csv", ".json", ".rst", ".log"}


class ExtractError(Exception):
    """Raised when a file can't be read or its format is unsupported."""

    def __init__(self, message: str, status: int = 415) -> None:
        super().__init__(message)
        self.message = message
        self.status = status


def extract_text(data: bytes, filename: str, *, char_limit: int = 24000) -> str:
    """Return extracted text from *data* (in-memory), capped to *char_limit*.

    Raises :class:`ExtractError` on an unsupported/failed format."""
    name = (filename or "").lower()
    suffix = "." + name.rsplit(".", 1)[-1] if "." in name else ""

    try:
        if suffix in _TEXT_SUFFIXES:
            text = data.decode("utf-8", errors="replace")
        elif suffix == ".pdf":
            text = _pdf(data)
        elif suffix == ".docx":
            text = _docx(data)
        elif suffix == ".pptx":
            text = _pptx(data)
        elif suffix == ".xlsx":
            text = _xlsx(data)
        elif suffix in {".odt", ".ods", ".odp"}:
            text = _odf(data)
        else:
            raise ExtractError(
                f"unsupported file type '{suffix or filename}'. "
                "Use txt/md/qmd/csv/json/pdf/docx/pptx/xlsx/odt, or paste the text.")
    except ExtractError:
        raise
    except Exception as exc:  # corrupted/encrypted file etc.
        raise ExtractError(f"couldn't read {filename}: {exc}") from exc

    text = (text or "").strip()
    if not text:
        raise ExtractError(f"no readable text found in {filename}.")
    return text[:char_limit]


def _pdf(data: bytes) -> str:
    from pypdf import PdfReader
    parts = []
    for page in PdfReader(io.BytesIO(data)).pages:
        parts.append(page.extract_text() or "")
    return "\n".join(parts)


def _docx(data: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _pptx(data: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(data))
    parts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _xlsx(data: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    out = []
    for ws in wb.worksheets:
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                out.append(" | ".join(cells))
    return "\n".join(out)


def _odf(data: bytes) -> str:
    """OpenDocument (odt/ods/odp) is a zip; pull text from content.xml."""
    with zipfile.ZipFile(io.BytesIO(data)) as zf:
        xml = zf.read("content.xml")
    # Naive text pull: concatenate all text-node content.
    root = ET.fromstring(xml)
    ns = "{urn:oasis:names:tc:opendocument:xmlns:text:1.0}"
    return "\n".join("".join(t.text or "" for t in p.iter(ns + "p")
                             if (t.text or "").strip()) for p in root.iter())
